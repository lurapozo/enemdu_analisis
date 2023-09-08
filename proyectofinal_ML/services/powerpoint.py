import pyautogui as py
import os
import time
from pptx.util import Inches,Pt

def gen_presentacion(prs, pos, titulo, contenido):

    if pos == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Numpy y Matrices"
        subtitle.text = "Grupo 4"
    
    elif pos == 1:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = titulo

        tf = body_shape.text_frame
        tf.text = contenido

    else:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes

        title_shape = shapes.title
        body_shape = shapes.placeholders[1]

        title_shape.text = titulo

        tf = body_shape.text_frame.paragraphs[0]
        tf.text = contenido
        tf.level = 2

    prs.save("clase.pptx")


def gen_img(prs, img_path):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    left = Inches(0.8)
    top = Inches(0.5)
    pic = slide.shapes.add_picture(img_path, left, top)
    prs.save("clase.pptx")
def init_presentetation(name):
    os.startfile(name)
    time.sleep(3)
    py.getWindowsWithTitle("PowerPoint")[0].maximize()
    py.press("f5")
    time.sleep(2)

def change_slide():
    py.press("right")

def end_presentation():
    py.press("esc")
    time.sleep(1)
